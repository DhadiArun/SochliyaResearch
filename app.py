"""
╔══════════════════════════════════════════════════════════════════════════════╗
║          MULTI-AGENT MARKET RESEARCH PLATFORM — COMPLETE APPLICATION        ║
║                                                                              ║
║  Architecture:  Dual-model research (Claude + Gemini) →                     ║
║                 Research Merger → LLM Council (3 agents) →                  ║
║                 Synthesis → Export (PDF / DOCX / PPTX)                      ║
║                                                                              ║
║  Stack:         Streamlit UI · Anthropic SDK · Google Generative AI         ║
║                 ReportLab · python-docx · python-pptx                       ║
║                                                                              ║
║  Setup:                                                                      ║
║    pip install streamlit anthropic google-generativeai                       ║
║               reportlab python-docx python-pptx                             ║
║                                                                              ║
║    Set environment variables (or use .env / Streamlit secrets):              ║
║    ANTHROPIC_API_KEY=sk-ant-...                                              ║
║    GOOGLE_API_KEY=AIza...                                                    ║
║                                                                              ║
║    Run:  streamlit run market_research_app.py                               ║
╚══════════════════════════════════════════════════════════════════════════════╝
"""

# ─── Standard library ─────────────────────────────────────────────────────────
import asyncio
import io
import json
import os
import re
import time
import datetime
from dataclasses import dataclass
from typing import Any

# ─── Third-party ──────────────────────────────────────────────────────────────
import streamlit as st

# ─── Lazy imports (allow app to load even before keys are entered) ─────────────
try:
    import anthropic
    ANTHROPIC_AVAILABLE = True
except ImportError:
    ANTHROPIC_AVAILABLE = False

try:
    import google.generativeai as genai
    GOOGLE_AVAILABLE = True
except ImportError:
    GOOGLE_AVAILABLE = False

try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table,
        TableStyle, HRFlowable, PageBreak
    )
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

try:
    from docx import Document
    from docx.shared import Pt, RGBColor, Cm
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt as PPTXPt
    from pptx.dml.color import RGBColor as PPTXColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 1 — CONFIGURATION & PROMPTS
# ══════════════════════════════════════════════════════════════════════════════

MODEL_CLAUDE_RESEARCH  = "claude-opus-4-6"          # Deep reasoning for research
MODEL_CLAUDE_MERGE     = "claude-haiku-4-5-20251001" # Fast + cheap for merge step
MODEL_CLAUDE_COUNCIL   = "claude-opus-4-6"           # Reasoning depth for critique
MODEL_CLAUDE_SYNTHESIS = "claude-opus-4-6"           # Highest quality for final report

RESEARCH_SYSTEM_PROMPT = """You are a senior market research analyst at a top-tier consultancy.
Produce a detailed, structured research report using ONLY the following markdown sections.
Every numeric claim MUST include a confidence qualifier (e.g., "estimated", "reported", "projected").
Flag unverifiable figures as [UNVERIFIED].
Never pad with generic observations — every sentence must add specific, actionable information.

## Market Size & Growth
(Total addressable market in USD billions, CAGR, key growth drivers, data methodology)

## Key Players
(Top 5–7 companies: market share estimate, key products, strategic positioning)

## Market Trends
(3–5 major trends with specific evidence and timeline)

## Risk Factors
(Top 3–5 risks: likelihood, potential impact, current mitigation landscape)

## Data Gaps & Limitations
(What critical information is missing or unverifiable)"""

MERGE_PROMPT_TEMPLATE = """You have two independent market research reports on: "{query}"

## Report A — Claude Analysis
{claude_output}

## Report B — Gemini Analysis
{gemini_output}

Task: Produce a single UNIFIED research context by:
1. Keeping ALL unique data points from both reports
2. Flagging contradictions as: [CONFLICT: Report A says X | Report B says Y]
3. Noting where both agree (higher confidence)
4. Assigning a confidence level (HIGH/MEDIUM/LOW) to each major claim
5. Outputting clean structured markdown (same sections as the input reports)

Output ONLY the merged markdown document."""

SYNTHESIS_SYSTEM_PROMPT = """You are a Chief Research Officer synthesising a multi-agent LLM council debate
into a final, boardroom-ready market research report.

Your responsibilities:
1. Weigh each council agent's critique by their stated domain expertise
2. Resolve [CONFLICT] flags by citing which evidence is stronger and why
3. REDUCE confidence scores where agents challenged claims without rebuttal
4. Ensure the final report is MORE accurate than any single agent could produce
5. Include an honest "council_dissent_log" of unresolved disagreements

Your output MUST be valid JSON — no markdown fences, no preamble, pure JSON only."""

FINAL_REPORT_SCHEMA = """{
  "executive_summary": "3-5 sentence boardroom summary",
  "market_overview": {
    "size_usd_billions": 0.0,
    "cagr_percent": 0.0,
    "forecast_year": 2029,
    "confidence_score": 75,
    "key_data_points": ["point 1", "point 2"]
  },
  "key_findings": [
    {"finding": "...", "confidence": 80, "source_models": ["Claude", "Gemini"]}
  ],
  "competitive_landscape": [
    {"player": "...", "position": "...", "threat_level": "HIGH|MED|LOW"}
  ],
  "risks": [
    {"risk": "...", "severity": "CRITICAL|HIGH|MED|LOW", "mitigation": "..."}
  ],
  "strategic_recommendations": [
    {"action": "...", "timeline": "...", "priority": 1}
  ],
  "council_dissent_log": [
    {"agent": "...", "unresolved_challenge": "..."}
  ],
  "overall_report_confidence": 75
}"""

@dataclass
class CouncilAgent:
    name:    str
    persona: str
    focus:   str
    color:   str  # For UI display

COUNCIL_AGENTS = [
    CouncilAgent(
        name="The Skeptic",
        persona="""You are a rigorous fact-checker and devil's advocate with 20 years of
research validation experience. Your SOLE purpose is to challenge every claim
in the research. You are NOT trying to be helpful — you are trying to find flaws.

Identify and call out:
- Unsupported assertions presented as facts
- Survivorship bias in competitive analysis
- Correlation/causation errors
- Overly optimistic market size projections
- Missing counter-evidence
- Methodological weaknesses in cited data

Output ONLY valid JSON matching this exact schema (no preamble, no fences):
{
  "agent": "The Skeptic",
  "verdict": "ACCEPT|CHALLENGE|REJECT",
  "confidence_adjustment": <integer from -30 to +5>,
  "key_findings": ["finding that survived scrutiny"],
  "challenges": ["specific challenged claim: reason"],
  "additions": ["critical missing context"],
  "risk_flags": ["specific risk the research underplayed"]
}""",
        focus="Factual accuracy & logical fallacies",
        color="#E24B4A",
    ),
    CouncilAgent(
        name="Financial Analyst",
        persona="""You are a CFO-level financial analyst who has reviewed thousands of
market research reports. You scrutinise ONLY the quantitative claims.

Your specific checks:
- Are market size figures internally consistent with growth rates?
- Do valuations align with comparable transactions?
- Are CAGR projections realistic given historical base rates?
- Is the TAM/SAM/SOM breakdown logical?
- Are revenue projections backed by a stated methodology?
- Do the numbers add up? (e.g., if 5 players share a $10B market, do their shares sum to ~100%?)

Output ONLY valid JSON matching this exact schema (no preamble, no fences):
{
  "agent": "Financial Analyst",
  "verdict": "ACCEPT|CHALLENGE|REJECT",
  "confidence_adjustment": <integer from -30 to +10>,
  "key_findings": ["validated financial finding"],
  "challenges": ["specific numeric claim challenged: issue"],
  "additions": ["missing financial context"],
  "risk_flags": ["financial risk underweighted"]
}""",
        focus="Quantitative claims & financial metrics",
        color="#BA7517",
    ),
    CouncilAgent(
        name="Industry Strategist",
        persona="""You are a McKinsey-level industry strategist who specialises in
competitive dynamics and market evolution. You evaluate strategic quality.

Your specific checks:
- Is the competitive landscape analysis complete? Who is missing?
- Are Porter's Five Forces dynamics correctly assessed?
- Are there disruptive threats (adjacent markets, technology shifts) not mentioned?
- Are the strategic recommendations actionable and specific?
- Is the market segmentation logical?
- What strategic opportunities does the research overlook?

Output ONLY valid JSON matching this exact schema (no preamble, no fences):
{
  "agent": "Industry Strategist",
  "verdict": "ACCEPT|CHALLENGE|REJECT",
  "confidence_adjustment": <integer from -20 to +10>,
  "key_findings": ["strategic insight confirmed"],
  "challenges": ["strategic gap or error identified"],
  "additions": ["competitive intelligence missing"],
  "risk_flags": ["strategic risk not addressed"]
}""",
        focus="Competitive dynamics & strategic quality",
        color="#1D9E75",
    ),
]


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 2 — LLM CLIENT HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def get_anthropic_client():
    key = st.session_state.get("anthropic_key") or os.environ.get("ANTHROPIC_API_KEY", "")
    if not key:
        raise ValueError("Anthropic API key not set. Enter it in the sidebar.")
    return anthropic.Anthropic(api_key=key)

def get_gemini_model():
    key = st.session_state.get("google_key") or os.environ.get("GOOGLE_API_KEY", "")
    if not key:
        raise ValueError("Google API key not set. Enter it in the sidebar.")
    genai.configure(api_key=key)
    return genai.GenerativeModel("gemini-1.5-pro")

def clean_json(raw: str) -> str:
    """Strip markdown code fences and stray text before/after JSON."""
    cleaned = re.sub(r"^```(?:json)?\s*", "", raw.strip(), flags=re.MULTILINE)
    cleaned = re.sub(r"```\s*$",           "", cleaned,     flags=re.MULTILINE)
    # Find the outermost { } in case there's preamble text
    start = cleaned.find("{")
    end   = cleaned.rfind("}") + 1
    if start >= 0 and end > start:
        return cleaned[start:end]
    return cleaned.strip()


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 3 — PIPELINE NODES
# ══════════════════════════════════════════════════════════════════════════════

def run_claude_research(query: str, status_cb) -> str:
    """Calls Claude for primary market research. Returns full markdown text."""
    status_cb("Calling Claude for primary research...")
    client = get_anthropic_client()
    full_text = ""
    with client.messages.stream(
        model=MODEL_CLAUDE_RESEARCH,
        max_tokens=4096,
        system=RESEARCH_SYSTEM_PROMPT,
        messages=[{"role": "user", "content": f"Research query: {query}"}],
    ) as stream:
        for text in stream.text_stream:
            full_text += text
    return full_text

def run_gemini_research(query: str, status_cb) -> str:
    """Calls Gemini 1.5 Pro for independent market research."""
    status_cb("Calling Gemini for independent research...")
    model = get_gemini_model()
    prompt = f"{RESEARCH_SYSTEM_PROMPT}\n\nResearch query: {query}"
    response = model.generate_content(
        prompt,
        generation_config=genai.GenerationConfig(
            max_output_tokens=4096,
            temperature=0.3,
        ),
    )
    return response.text

def run_merge(query: str, claude_out: str, gemini_out: str, status_cb) -> str:
    """Merges both research outputs using Claude Haiku (fast + cheap)."""
    status_cb("Merging research outputs and detecting conflicts...")
    client = get_anthropic_client()
    prompt = MERGE_PROMPT_TEMPLATE.format(
        query=query,
        claude_output=claude_out[:6000],   # Truncate to stay within context
        gemini_output=gemini_out[:6000],
    )
    response = client.messages.create(
        model=MODEL_CLAUDE_MERGE,
        max_tokens=6000,
        messages=[{"role": "user", "content": prompt}],
    )
    return response.content[0].text

def run_council_agent(agent: CouncilAgent, merged: str, query: str, status_cb) -> dict:
    """Runs one council agent. Returns parsed JSON review dict."""
    status_cb(f"Council: {agent.name} reviewing...")
    client = get_anthropic_client()
    prompt = f"""Research Query: {query}

=== UNIFIED RESEARCH CONTEXT ===
{merged[:8000]}
=================================

Your focus area: {agent.focus}

Produce your council review now. Output ONLY the JSON object — no preamble, no markdown."""

    response = client.messages.create(
        model=MODEL_CLAUDE_COUNCIL,
        max_tokens=2048,
        system=agent.persona,
        messages=[{"role": "user", "content": prompt}],
    )
    raw = response.content[0].text
    try:
        return json.loads(clean_json(raw))
    except json.JSONDecodeError:
        # Graceful degradation — return partial structured data
        return {
            "agent":               agent.name,
            "verdict":             "CHALLENGE",
            "confidence_adjustment": 0,
            "key_findings":        [],
            "challenges":          [f"Agent returned unparseable output: {raw[:200]}"],
            "additions":           [],
            "risk_flags":          [],
        }

def run_synthesis(query: str, merged: str, reviews: list[dict], status_cb) -> dict:
    """Synthesis agent compiles council debate into final structured report."""
    status_cb("Synthesis agent compiling final report...")
    client = get_anthropic_client()
    prompt = f"""Original Query: {query}

=== MERGED RESEARCH BASE ===
{merged[:6000]}

=== COUNCIL DEBATE ({len(reviews)} agents) ===
{json.dumps(reviews, indent=2)}

=== TARGET OUTPUT SCHEMA ===
{FINAL_REPORT_SCHEMA}

Synthesise the above into a final report matching the schema exactly.
Resolve all [CONFLICT] flags. Reduce confidence where agents challenged claims.
Output ONLY valid JSON — no markdown, no preamble."""

    response = client.messages.create(
        model=MODEL_CLAUDE_SYNTHESIS,
        max_tokens=8000,
        system=SYNTHESIS_SYSTEM_PROMPT,
        messages=[{"role": "user", "content": prompt}],
    )
    raw = response.content[0].text
    try:
        return json.loads(clean_json(raw))
    except json.JSONDecodeError as e:
        st.error(f"Synthesis JSON parse error: {e}. Returning raw text.")
        return {"executive_summary": raw, "parse_error": True}

def run_full_pipeline(query: str, status_placeholder) -> dict:
    """
    Orchestrates the full multi-agent pipeline.
    Returns a state dict containing all intermediate + final outputs.
    """
    results: dict[str, Any] = {"query": query}

    def status(msg: str):
        status_placeholder.info(f"⚙️ {msg}")

    # ── Node 1: Parallel Research ──────────────────────────────────────────
    # Note: Streamlit is synchronous; we run models sequentially.
    # In production FastAPI, replace with asyncio.gather for true parallelism.
    status("Starting dual-model research phase...")
    t0 = time.time()

    try:
        claude_out = run_claude_research(query, status)
        results["claude_output"] = claude_out
        status(f"Claude research complete ({len(claude_out):,} chars).")
    except Exception as e:
        results["claude_output"] = f"[Claude error: {e}]"
        st.warning(f"Claude research failed: {e}")

    try:
        gemini_out = run_gemini_research(query, status)
        results["gemini_output"] = gemini_out
        status(f"Gemini research complete ({len(gemini_out):,} chars).")
    except Exception as e:
        results["gemini_output"] = f"[Gemini error: {e}]"
        st.warning(f"Gemini research failed: {e}")

    results["research_time_s"] = round(time.time() - t0, 1)

    # ── Node 2: Merge ──────────────────────────────────────────────────────
    t1 = time.time()
    try:
        merged = run_merge(query, results["claude_output"], results["gemini_output"], status)
        results["merged_context"] = merged
        conflict_count = merged.count("[CONFLICT:")
        status(f"Merge complete. Found {conflict_count} conflict(s) between models.")
    except Exception as e:
        results["merged_context"] = results.get("claude_output", "")
        st.warning(f"Merge failed, using Claude output only: {e}")
    results["merge_time_s"] = round(time.time() - t1, 1)

    # ── Node 3: LLM Council ────────────────────────────────────────────────
    t2 = time.time()
    reviews = []
    for agent in COUNCIL_AGENTS:
        try:
            review = run_council_agent(agent, results["merged_context"], query, status)
            reviews.append(review)
            verdict = review.get("verdict", "?")
            adj     = review.get("confidence_adjustment", 0)
            adj_str = f"+{adj}" if adj > 0 else str(adj)
            status(f"{agent.name}: verdict={verdict}, confidence_adj={adj_str}")
        except Exception as e:
            st.warning(f"Council agent {agent.name} failed: {e}")
    results["council_reviews"] = reviews
    results["council_time_s"] = round(time.time() - t2, 1)

    # ── Node 4: Synthesis ──────────────────────────────────────────────────
    t3 = time.time()
    try:
        final = run_synthesis(query, results["merged_context"], reviews, status)
        results["final_report"] = final
        conf = final.get("overall_report_confidence", "?")
        status(f"Synthesis complete. Overall confidence: {conf}%")
    except Exception as e:
        results["final_report"] = {}
        st.error(f"Synthesis failed: {e}")
    results["synthesis_time_s"] = round(time.time() - t3, 1)

    results["total_time_s"] = round(time.time() - t0, 1)
    status_placeholder.success(f"✅ Pipeline complete in {results['total_time_s']}s")
    return results


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 4 — PDF EXPORT
# ══════════════════════════════════════════════════════════════════════════════

def export_pdf(report: dict, query: str) -> bytes:
    if not REPORTLAB_AVAILABLE:
        raise ImportError("reportlab not installed.")

    NAVY   = colors.HexColor("#0F2040")
    TEAL   = colors.HexColor("#1D9E75")
    AMBER  = colors.HexColor("#BA7517")
    LIGHT  = colors.HexColor("#F1EFE8")
    BORDER = colors.HexColor("#B4B2A9")
    MUTED  = colors.HexColor("#5F5E5A")
    RED    = colors.HexColor("#E24B4A")
    CRIT   = colors.HexColor("#7F1515")
    WHITE  = colors.white
    BLACK  = colors.HexColor("#2C2C2A")

    SEV_COLOR = {"LOW": TEAL, "MED": AMBER, "HIGH": RED, "CRITICAL": CRIT}
    PRI_COLOR = {1: TEAL, 2: AMBER, 3: BORDER}

    buf = io.BytesIO()
    base = getSampleStyleSheet()

    def sty(name, **kw):
        return ParagraphStyle(name, parent=base["Normal"], **kw)

    h1   = sty("H1",   fontName="Helvetica-Bold", fontSize=18, textColor=NAVY,   spaceBefore=16, spaceAfter=6,  leading=24)
    h2   = sty("H2",   fontName="Helvetica-Bold", fontSize=13, textColor=TEAL,   spaceBefore=12, spaceAfter=4,  leading=18)
    body = sty("Body", fontName="Helvetica",       fontSize=10, textColor=BLACK,  leading=16,     spaceAfter=4,  alignment=TA_JUSTIFY)
    blt  = sty("Blt",  fontName="Helvetica",       fontSize=10, textColor=BLACK,  leading=16,     leftIndent=14, spaceAfter=3)
    smll = sty("Sm",   fontName="Helvetica",       fontSize=8,  textColor=MUTED,  leading=12,     spaceAfter=2)
    thdr = sty("TH",   fontName="Helvetica-Bold",  fontSize=9,  textColor=WHITE,  leading=13)
    tcll = sty("TC",   fontName="Helvetica",       fontSize=9,  textColor=BLACK,  leading=13)
    cvrt = sty("Cvr",  fontName="Helvetica-Bold",  fontSize=28, textColor=WHITE,  alignment=TA_CENTER, leading=36)
    cvrs = sty("CvrS", fontName="Helvetica",       fontSize=13, textColor=colors.HexColor("#9FE1CB"), alignment=TA_CENTER, leading=20)
    cvrm = sty("CvrM", fontName="Helvetica",       fontSize=9,  textColor=BORDER, alignment=TA_CENTER, leading=14)

    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=2*cm, rightMargin=2*cm,
        topMargin=2.2*cm, bottomMargin=2.5*cm,
    )

    story = []
    now   = datetime.datetime.now().strftime("%B %d, %Y")

    def sp(n=6):  return Spacer(1, n)
    def hr_line(): return HRFlowable(width="100%", thickness=1, color=TEAL, spaceAfter=8)

    def tbl_of(rows, hdrs=None, cws=None):
        data = []
        if hdrs:
            data.append([Paragraph(h, thdr) for h in hdrs])
        for row in rows:
            data.append([Paragraph(str(c), tcll) for c in row])
        t = Table(data, colWidths=cws, repeatRows=1 if hdrs else 0)
        ts = [
            ("GRID", (0,0), (-1,-1), 0.25, BORDER),
            ("TOPPADDING", (0,0), (-1,-1), 5), ("BOTTOMPADDING", (0,0), (-1,-1), 5),
            ("LEFTPADDING", (0,0), (-1,-1), 7), ("RIGHTPADDING", (0,0), (-1,-1), 7),
            ("VALIGN", (0,0), (-1,-1), "TOP"),
        ]
        if hdrs:
            ts += [("BACKGROUND", (0,0), (-1,0), NAVY),
                   ("ROWBACKGROUNDS", (0,1), (-1,-1), [WHITE, LIGHT])]
        t.setStyle(TableStyle(ts))
        return t

    # ── Cover ──────────────────────────────────────────────────────────────────
    cover_rows = [
        [sp(50)], [Paragraph("MARKET RESEARCH REPORT", cvrt)],
        [sp(8)],  [Paragraph("LLM Council Validated Analysis", cvrs)],
        [sp(16)], [HRFlowable(width="60%", thickness=1, color=TEAL, spaceAfter=12)],
        [Paragraph(f"Query: {query}", cvrm)],
        [sp(4)],  [Paragraph(f"Generated: {now}  ·  Confidence: {report.get('overall_report_confidence','N/A')}%", cvrm)],
    ]
    cover = Table([[r[0]] for r in cover_rows], colWidths=["100%"])
    cover.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,-1), NAVY),
        ("LEFTPADDING", (0,0), (-1,-1), 30), ("RIGHTPADDING", (0,0), (-1,-1), 30),
        ("TOPPADDING", (0,0), (-1,-1), 0), ("BOTTOMPADDING", (0,0), (-1,-1), 0),
    ]))
    story.append(cover)
    story.append(PageBreak())

    # ── Executive Summary ──────────────────────────────────────────────────────
    story.append(Paragraph("Executive Summary", h1))
    story.append(hr_line())
    story.append(Paragraph(report.get("executive_summary", "N/A"), body))
    story.append(sp(12))

    # ── Market Overview ────────────────────────────────────────────────────────
    mo = report.get("market_overview", {})
    story.append(Paragraph("Market Overview", h1))
    story.append(hr_line())
    mo_data = [
        ["Metric", "Value", "Confidence"],
        ["Market Size (USD B)",  f"${mo.get('size_usd_billions','N/A')}B", f"{mo.get('confidence_score','N/A')}%"],
        ["CAGR",                 f"{mo.get('cagr_percent','N/A')}%",        ""],
        ["Forecast Year",        str(mo.get("forecast_year","N/A")),         ""],
    ]
    story.append(tbl_of(mo_data[1:], mo_data[0], [6.5*cm, 5*cm, 4.5*cm]))
    story.append(sp(8))
    for pt in mo.get("key_data_points", []):
        story.append(Paragraph(f"\u2022  {pt}", blt))
    story.append(sp(12))

    # ── Key Findings ───────────────────────────────────────────────────────────
    story.append(Paragraph("Key Findings", h1))
    story.append(hr_line())
    for i, f in enumerate(report.get("key_findings", []), 1):
        conf = f.get("confidence", 0)
        col  = "#1D9E75" if conf >= 70 else "#BA7517" if conf >= 40 else "#E24B4A"
        story.append(Paragraph(
            f"<b>{i}.</b> {f.get('finding','')} "
            f"<font color='{col}' size='8'>[{conf}% confidence]</font>",
            body))
        story.append(sp(3))
    story.append(sp(8))

    # ── Competitive Landscape ──────────────────────────────────────────────────
    story.append(Paragraph("Competitive Landscape", h1))
    story.append(hr_line())
    comp = report.get("competitive_landscape", [])
    if comp:
        comp_rows = [[c.get("player",""), c.get("position",""), c.get("threat_level","")] for c in comp]
        ct = tbl_of(comp_rows, ["Player", "Strategic Position", "Threat Level"], [4*cm, 9*cm, 3*cm])
        # Colour threat level cells
        for ri, c in enumerate(comp, start=1):
            tl  = c.get("threat_level","")
            clr = {"HIGH": RED, "MED": AMBER, "LOW": TEAL}.get(tl, BORDER)
            ct.setStyle(TableStyle([
                ("BACKGROUND", (2, ri), (2, ri), clr),
                ("TEXTCOLOR",  (2, ri), (2, ri), WHITE),
            ]))
        story.append(ct)
    story.append(sp(12))

    # ── Risks ──────────────────────────────────────────────────────────────────
    story.append(PageBreak())
    story.append(Paragraph("Risk Register", h1))
    story.append(hr_line())
    risks = report.get("risks", [])
    if risks:
        risk_rows = [[r.get("risk",""), r.get("severity",""), r.get("mitigation","")] for r in risks]
        rt = tbl_of(risk_rows, ["Risk", "Severity", "Mitigation Strategy"], [5.5*cm, 2.5*cm, 8*cm])
        for ri, r in enumerate(risks, start=1):
            sc = SEV_COLOR.get(r.get("severity",""), BORDER)
            rt.setStyle(TableStyle([("BACKGROUND",(1,ri),(1,ri),sc), ("TEXTCOLOR",(1,ri),(1,ri),WHITE)]))
        story.append(rt)
    story.append(sp(12))

    # ── Strategic Recommendations ──────────────────────────────────────────────
    story.append(Paragraph("Strategic Recommendations", h1))
    story.append(hr_line())
    recs = sorted(report.get("strategic_recommendations", []), key=lambda x: x.get("priority", 9))
    for rec in recs:
        p    = rec.get("priority", 3)
        pclr = {1:"#1D9E75", 2:"#BA7517", 3:"#888780"}.get(p, "#888780")
        story.append(Paragraph(
            f"<font color='{pclr}'><b>[P{p}]</b></font>  <b>{rec.get('action','')}</b>"
            f"  <font color='#5F5E5A' size='9'>— {rec.get('timeline','')}</font>",
            body))
        story.append(sp(5))
    story.append(sp(10))

    # ── Council Dissent Log ────────────────────────────────────────────────────
    dissent = report.get("council_dissent_log", [])
    if dissent:
        story.append(Paragraph("Council Dissent Log", h1))
        story.append(hr_line())
        for d in dissent:
            story.append(Paragraph(
                f"<b>{d.get('agent','?')}</b>: {d.get('unresolved_challenge','')}", blt))
            story.append(sp(3))

    story.append(sp(16))
    story.append(HRFlowable(width="100%", thickness=0.5, color=BORDER))
    story.append(Paragraph(
        f"Multi-Agent Market Research Platform · Generated {now} · "
        f"Overall Confidence: {report.get('overall_report_confidence','N/A')}%",
        smll))

    doc.build(story)
    return buf.getvalue()


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 5 — DOCX EXPORT
# ══════════════════════════════════════════════════════════════════════════════

def export_docx(report: dict, query: str) -> bytes:
    if not DOCX_AVAILABLE:
        raise ImportError("python-docx not installed.")

    NAVY_RGB  = RGBColor(0x0F, 0x20, 0x40)
    TEAL_RGB  = RGBColor(0x1D, 0x9E, 0x75)
    AMBER_RGB = RGBColor(0xBA, 0x75, 0x17)
    RED_RGB   = RGBColor(0xE2, 0x4B, 0x4A)
    CRIT_RGB  = RGBColor(0x7F, 0x15, 0x15)
    WHITE_RGB = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_HEX = "F1EFE8"

    SEV_HEX = {"LOW":"1D9E75","MED":"BA7517","HIGH":"E24B4A","CRITICAL":"7F1515"}
    TH_HEX  = {"HIGH":"E24B4A","MED":"BA7517","LOW":"1D9E75"}

    def set_cell_bg(cell, hex_c: str):
        tc   = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd  = OxmlElement("w:shd")
        shd.set(qn("w:val"),  "clear")
        shd.set(qn("w:color"),"auto")
        shd.set(qn("w:fill"), hex_c)
        tcPr.append(shd)

    def hdr(doc, text, level, color=NAVY_RGB, size=None):
        h = doc.add_heading(text, level=level)
        r = h.runs[0] if h.runs else h.add_run(text)
        r.font.color.rgb = color
        if size: r.font.size = Pt(size)
        return h

    def tbl(doc, header_row, data_rows, col_widths_cm, hdr_bg="0F2040"):
        n_cols = len(header_row)
        t = doc.add_table(rows=1 + len(data_rows), cols=n_cols)
        t.style = "Table Grid"
        for i, h_text in enumerate(header_row):
            cell = t.rows[0].cells[i]
            cell.width = Cm(col_widths_cm[i])
            p    = cell.paragraphs[0]
            run  = p.add_run(h_text)
            run.font.bold = True
            run.font.size = Pt(9)
            run.font.color.rgb = WHITE_RGB
            set_cell_bg(cell, hdr_bg)
        for ri, row in enumerate(data_rows, start=1):
            bg = LIGHT_HEX if ri % 2 == 0 else "FFFFFF"
            for ci, val in enumerate(row):
                cell = t.rows[ri].cells[ci]
                cell.width = Cm(col_widths_cm[ci])
                cell.paragraphs[0].add_run(str(val)).font.size = Pt(9)
                set_cell_bg(cell, bg)
        return t

    doc = Document()
    for sec in doc.sections:
        sec.top_margin    = Cm(2.5)
        sec.bottom_margin = Cm(2.5)
        sec.left_margin   = Cm(2.8)
        sec.right_margin  = Cm(2.8)

    now = datetime.datetime.now().strftime("%B %d, %Y")

    # Title block
    hdr(doc, "Market Research Report", 1, NAVY_RGB, 22)
    meta = doc.add_paragraph(f"Query: {query}")
    meta.runs[0].font.size  = Pt(9)
    meta.runs[0].font.color.rgb = RGBColor(0x88,0x87,0x80)
    meta2 = doc.add_paragraph(
        f"Generated: {now}  ·  Council Confidence: {report.get('overall_report_confidence','N/A')}%")
    meta2.runs[0].font.size = Pt(9)
    meta2.runs[0].font.color.rgb = RGBColor(0x88,0x87,0x80)
    doc.add_paragraph()

    # Executive Summary
    hdr(doc, "Executive Summary", 2, TEAL_RGB, 14)
    doc.add_paragraph(report.get("executive_summary","N/A"))
    doc.add_paragraph()

    # Market Overview
    hdr(doc, "Market Overview", 2, TEAL_RGB, 14)
    mo = report.get("market_overview",{})
    tbl(doc,
        ["Metric","Value","Confidence"],
        [
            ("Market Size (USD B)", f"${mo.get('size_usd_billions','N/A')}B", f"{mo.get('confidence_score','N/A')}%"),
            ("CAGR",               f"{mo.get('cagr_percent','N/A')}%",        ""),
            ("Forecast Year",      str(mo.get("forecast_year","N/A")),         ""),
        ],
        [6, 5, 4])
    doc.add_paragraph()
    for pt in mo.get("key_data_points",[]):
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(pt).font.size = Pt(10)
    doc.add_paragraph()

    # Key Findings
    hdr(doc, "Key Findings", 2, TEAL_RGB, 14)
    for i, f in enumerate(report.get("key_findings",[]), 1):
        conf = f.get("confidence",0)
        p    = doc.add_paragraph(style="List Number")
        r1   = p.add_run(f.get("finding",""))
        r1.font.size = Pt(10)
        r2   = p.add_run(f"  [{conf}% confidence]")
        r2.font.size = Pt(8)
        r2.font.color.rgb = TEAL_RGB if conf>=70 else AMBER_RGB
    doc.add_paragraph()

    # Competitive Landscape
    hdr(doc, "Competitive Landscape", 2, TEAL_RGB, 14)
    comp = report.get("competitive_landscape",[])
    if comp:
        ct = tbl(doc,
            ["Player","Strategic Position","Threat Level"],
            [(c.get("player",""), c.get("position",""), c.get("threat_level","")) for c in comp],
            [4,9,3])
        for ri, c in enumerate(comp, start=1):
            tl  = c.get("threat_level","")
            hex = TH_HEX.get(tl,"B4B2A9")
            cell = ct.rows[ri].cells[2]
            set_cell_bg(cell, hex)
            cell.paragraphs[0].runs[0].font.color.rgb = WHITE_RGB
    doc.add_paragraph()

    # Risk Register
    hdr(doc, "Risk Register", 2, TEAL_RGB, 14)
    risks = report.get("risks",[])
    if risks:
        rt = tbl(doc,
            ["Risk","Severity","Mitigation Strategy"],
            [(r.get("risk",""), r.get("severity",""), r.get("mitigation","")) for r in risks],
            [5.5, 2.5, 8])
        for ri, r in enumerate(risks, start=1):
            hex = SEV_HEX.get(r.get("severity",""),"B4B2A9")
            cell = rt.rows[ri].cells[1]
            set_cell_bg(cell, hex)
            cell.paragraphs[0].runs[0].font.color.rgb = WHITE_RGB
    doc.add_paragraph()

    # Strategic Recommendations
    hdr(doc, "Strategic Recommendations", 2, TEAL_RGB, 14)
    for rec in sorted(report.get("strategic_recommendations",[]), key=lambda x: x.get("priority",9)):
        p  = doc.add_paragraph(style="List Number")
        r1 = p.add_run(f"[P{rec.get('priority','')}] {rec.get('action','')}")
        r1.font.bold = True; r1.font.size = Pt(10)
        r2 = p.add_run(f" — {rec.get('timeline','')}")
        r2.font.size = Pt(9); r2.font.color.rgb = AMBER_RGB
    doc.add_paragraph()

    # Council Dissent Log
    dissent = report.get("council_dissent_log",[])
    if dissent:
        hdr(doc, "Council Dissent Log", 2, TEAL_RGB, 14)
        for d in dissent:
            p = doc.add_paragraph(style="List Bullet")
            r1= p.add_run(f"{d.get('agent','?')}: ")
            r1.font.bold = True; r1.font.size = Pt(10)
            p.add_run(d.get("unresolved_challenge","")).font.size = Pt(10)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 6 — PPTX EXPORT
# ══════════════════════════════════════════════════════════════════════════════

def export_pptx(report: dict, query: str) -> bytes:
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx not installed.")

    W, H = Inches(13.33), Inches(7.5)

    NAVY  = PPTXColor(0x0F, 0x20, 0x40)
    TEAL  = PPTXColor(0x1D, 0x9E, 0x75)
    AMBER = PPTXColor(0xBA, 0x75, 0x17)
    RED   = PPTXColor(0xE2, 0x4B, 0x4A)
    CRIT  = PPTXColor(0x7F, 0x15, 0x15)
    WHITE = PPTXColor(0xFF, 0xFF, 0xFF)
    LIGHT = PPTXColor(0xF1, 0xEF, 0xE8)
    MUTED = PPTXColor(0x5F, 0x5E, 0x5A)

    SEV_CLR = {"LOW":TEAL,"MED":AMBER,"HIGH":RED,"CRITICAL":CRIT}
    PRI_CLR = {1:TEAL, 2:AMBER, 3:PPTXColor(0xB4,0xB2,0xA9)}

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]  # Blank layout

    def add_slide():
        return prs.slides.add_slide(blank)

    def rect(slide, x, y, w, h, fill_color, line=False):
        from pptx.util import Emu
        shp = slide.shapes.add_shape(1, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
        if not line:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = fill_color
        return shp

    def txbox(slide, text, x, y, w, h, size=PPTXPt(11), bold=False,
              color=None, align=PP_ALIGN.LEFT, wrap=True):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = wrap
        p  = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size  = size
        run.font.bold  = bold
        run.font.color.rgb = color or PPTXColor(0x0F,0x20,0x40)
        return tb

    def title_bar(slide, title: str):
        bar = rect(slide, 0, 0, W, Inches(0.85), NAVY)
        tf  = bar.text_frame
        tf.text = title
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
        r = tf.paragraphs[0].runs[0]
        r.font.size = PPTXPt(22); r.font.bold = True; r.font.color.rgb = WHITE
        tf.margin_left = Inches(0.4); tf.margin_top = Inches(0.13)

    now = datetime.datetime.now().strftime("%B %d, %Y")

    # ── Slide 1: Cover ─────────────────────────────────────────────────────────
    s1 = add_slide()
    rect(s1, 0, 0, W, H, NAVY)
    txbox(s1, "MARKET RESEARCH REPORT",
          Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.4),
          size=PPTXPt(38), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s1, "LLM Council Validated Analysis",
          Inches(0.8), Inches(3.2), Inches(11.5), Inches(0.7),
          size=PPTXPt(18), color=PPTXColor(0x9F,0xE1,0xCB), align=PP_ALIGN.CENTER)
    txbox(s1, query,
          Inches(1.5), Inches(4.1), Inches(10), Inches(0.8),
          size=PPTXPt(14), color=PPTXColor(0xB4,0xB2,0xA9), align=PP_ALIGN.CENTER)
    txbox(s1, f"Generated: {now}  ·  Confidence: {report.get('overall_report_confidence','N/A')}%",
          Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5),
          size=PPTXPt(11), color=PPTXColor(0x88,0x87,0x80), align=PP_ALIGN.CENTER)

    # ── Slide 2: Executive Summary ─────────────────────────────────────────────
    s2 = add_slide()
    title_bar(s2, "Executive Summary")
    txbox(s2, report.get("executive_summary","N/A"),
          Inches(0.5), Inches(1.0), Inches(12.3), Inches(3.2),
          size=PPTXPt(13), color=NAVY)
    mo = report.get("market_overview",{})
    metrics = [
        ("Market Size", f"${mo.get('size_usd_billions','N/A')}B"),
        ("CAGR",        f"{mo.get('cagr_percent','N/A')}%"),
        ("Confidence",  f"{mo.get('overall_report_confidence',report.get('overall_report_confidence','N/A'))}%"),
    ]
    for i, (label, val) in enumerate(metrics):
        x = Inches(0.5 + i * 4.2)
        card = rect(s2, x, Inches(4.4), Inches(3.9), Inches(2.7), LIGHT)
        card.line.color.rgb = TEAL
        txbox(s2, label, x+Inches(0.15), Inches(4.6), Inches(3.6), Inches(0.6),
              size=PPTXPt(11), color=TEAL, bold=True)
        txbox(s2, val,   x+Inches(0.15), Inches(5.2), Inches(3.6), Inches(1.5),
              size=PPTXPt(32), bold=True, color=NAVY)

    # ── Slide 3: Key Findings ──────────────────────────────────────────────────
    s3 = add_slide()
    title_bar(s3, "Key Findings")
    findings = report.get("key_findings",[])[:6]
    for i, f in enumerate(findings):
        y    = Inches(1.0 + i * 1.06)
        conf = f.get("confidence",0)
        dot  = rect(s3, Inches(0.3), y+Inches(0.08), Inches(0.35), Inches(0.35),
                    TEAL if conf>=70 else AMBER)
        txbox(s3, f.get("finding",""),
              Inches(0.8), y, Inches(11), Inches(0.9), size=PPTXPt(11), color=NAVY)
        txbox(s3, f"{conf}%",
              Inches(12.0), y, Inches(1.0), Inches(0.5),
              size=PPTXPt(10), color=TEAL if conf>=70 else AMBER,
              bold=True, align=PP_ALIGN.RIGHT)

    # ── Slide 4: Risk Register ─────────────────────────────────────────────────
    s4 = add_slide()
    title_bar(s4, "Risk Register")
    risks = report.get("risks",[])[:5]
    for i, risk in enumerate(risks):
        y   = Inches(0.95 + i * 1.26)
        sev = risk.get("severity","")
        sc  = SEV_CLR.get(sev, AMBER)
        badge = rect(s4, Inches(0.3), y, Inches(1.3), Inches(0.52), sc)
        tf = badge.text_frame; tf.text = sev
        tf.paragraphs[0].runs[0].font.color.rgb = WHITE
        tf.paragraphs[0].runs[0].font.size = PPTXPt(9)
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.margin_top = Inches(0.06)
        txbox(s4, risk.get("risk",""),
              Inches(1.8), y, Inches(5.5), Inches(1.1), size=PPTXPt(11), bold=True, color=NAVY)
        txbox(s4, risk.get("mitigation",""),
              Inches(7.5), y, Inches(5.6), Inches(1.1), size=PPTXPt(10), color=MUTED)

    # ── Slide 5: Strategic Recommendations ────────────────────────────────────
    s5 = add_slide()
    title_bar(s5, "Strategic Recommendations")
    recs = sorted(report.get("strategic_recommendations",[]), key=lambda x: x.get("priority",9))[:5]
    for i, rec in enumerate(recs):
        y    = Inches(0.95 + i * 1.26)
        p    = rec.get("priority",3)
        pclr = PRI_CLR.get(p, PPTXColor(0xB4,0xB2,0xA9))
        num  = rect(s5, Inches(0.3), y, Inches(0.5), Inches(0.5), pclr)
        tf   = num.text_frame; tf.text = str(p)
        tf.paragraphs[0].runs[0].font.color.rgb = WHITE
        tf.paragraphs[0].runs[0].font.size = PPTXPt(14)
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.margin_top = Inches(0.06)
        txbox(s5, rec.get("action",""),
              Inches(1.0), y, Inches(10.5), Inches(0.7), size=PPTXPt(12), bold=True, color=NAVY)
        txbox(s5, rec.get("timeline",""),
              Inches(1.0), y+Inches(0.65), Inches(10.5), Inches(0.45),
              size=PPTXPt(10), color=AMBER)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 7 — STREAMLIT UI
# ══════════════════════════════════════════════════════════════════════════════

st.set_page_config(
    page_title="Multi-Agent Market Research",
    page_icon="🔬",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ── Custom CSS ─────────────────────────────────────────────────────────────────
st.markdown("""
<style>
.main-header {
    background: linear-gradient(135deg, #0F2040 0%, #1a3a6e 100%);
    padding: 2rem 2.5rem; border-radius: 12px; margin-bottom: 1.5rem;
}
.main-header h1 { color: white; font-size: 2rem; margin: 0; font-weight: 700; }
.main-header p  { color: #9FE1CB; margin: 0.3rem 0 0; font-size: 1rem; }
.agent-card {
    border-radius: 10px; padding: 1rem 1.2rem; margin: 0.5rem 0;
    border-left: 4px solid;
}
.verdict-badge {
    display: inline-block; padding: 2px 10px; border-radius: 20px;
    font-size: 0.75rem; font-weight: 700; color: white;
}
.metric-card {
    background: #F1EFE8; border-radius: 10px; padding: 1rem;
    border-left: 4px solid #1D9E75; text-align: center;
}
.metric-card .value { font-size: 2rem; font-weight: 700; color: #0F2040; }
.metric-card .label { font-size: 0.8rem; color: #5F5E5A; text-transform: uppercase; }
.section-header {
    border-bottom: 2px solid #1D9E75; padding-bottom: 0.3rem;
    margin: 1.5rem 0 0.8rem; color: #0F2040; font-weight: 700;
}
.finding-row {
    padding: 0.6rem 1rem; border-radius: 8px; margin: 0.3rem 0;
    background: #F8F8F6; border-left: 3px solid;
}
.risk-high   { background: #FCEBEB; border-left: 4px solid #E24B4A; }
.risk-med    { background: #FAEEDA; border-left: 4px solid #BA7517; }
.risk-low    { background: #EAF3DE; border-left: 4px solid #1D9E75; }
.risk-crit   { background: #F7C1C1; border-left: 4px solid #7F1515; }
stButton button { border-radius: 8px; font-weight: 600; }
</style>
""", unsafe_allow_html=True)

# ── Sidebar — API Keys & Settings ──────────────────────────────────────────────
with st.sidebar:
    st.markdown("## ⚙️ Configuration")
    st.markdown("---")
    st.markdown("**API Keys**")
    anthropic_key = st.text_input(
        "Anthropic API Key",
        type="password",
        value=st.session_state.get("anthropic_key",""),
        placeholder="sk-ant-...",
        help="Used for Claude research, merge, council, and synthesis.",
    )
    google_key = st.text_input(
        "Google AI API Key",
        type="password",
        value=st.session_state.get("google_key",""),
        placeholder="AIza...",
        help="Used for Gemini 1.5 Pro independent research.",
    )
    if anthropic_key: st.session_state["anthropic_key"] = anthropic_key
    if google_key:    st.session_state["google_key"]    = google_key

    st.markdown("---")
    st.markdown("**Council Settings**")
    use_skeptic    = st.checkbox("Enable The Skeptic",          value=True)
    use_analyst    = st.checkbox("Enable Financial Analyst",    value=True)
    use_strategist = st.checkbox("Enable Industry Strategist",  value=True)

    st.markdown("---")
    st.markdown("**Library Status**")
    def lib_status(available, name):
        icon = "✅" if available else "❌"
        st.markdown(f"{icon} {name}")
    lib_status(ANTHROPIC_AVAILABLE, "anthropic")
    lib_status(GOOGLE_AVAILABLE,    "google-generativeai")
    lib_status(REPORTLAB_AVAILABLE, "reportlab (PDF)")
    lib_status(DOCX_AVAILABLE,      "python-docx (DOCX)")
    lib_status(PPTX_AVAILABLE,      "python-pptx (PPTX)")

    st.markdown("---")
    st.markdown("""
**Pipeline Flow:**
1. Claude + Gemini (parallel)
2. Merge + conflict detection
3. LLM Council (3 agents)
4. Synthesis agent
5. Export (PDF/DOCX/PPTX)
""")

# ── Main header ────────────────────────────────────────────────────────────────
st.markdown("""
<div class="main-header">
  <h1>🔬 Multi-Agent Market Research Platform</h1>
  <p>Dual-model research · LLM Council vetting · Confidence-scored synthesis · Professional export</p>
</div>
""", unsafe_allow_html=True)

# ── Query input ────────────────────────────────────────────────────────────────
col_q, col_btn = st.columns([5, 1])
with col_q:
    query = st.text_input(
        "Research Query",
        placeholder='e.g. "Global AI chipset market 2024–2029 — key players, growth drivers, risks"',
        label_visibility="collapsed",
    )
with col_btn:
    run_btn = st.button("▶ Run Research", type="primary", use_container_width=True)

# ── Run pipeline ───────────────────────────────────────────────────────────────
if run_btn:
    if not query.strip():
        st.warning("Please enter a research query.")
    elif not anthropic_key:
        st.error("Anthropic API key is required.")
    else:
        global COUNCIL_AGENTS
        # Filter active council agents
        active_agents = []
        if use_skeptic:    active_agents.append(COUNCIL_AGENTS[0])
        if use_analyst:    active_agents.append(COUNCIL_AGENTS[1])
        if use_strategist: active_agents.append(COUNCIL_AGENTS[2])

        # Temporarily override global list for this run
        COUNCIL_AGENTS = active_agents

        status_placeholder = st.empty()

        with st.spinner("Running multi-agent pipeline..."):
            results = run_full_pipeline(query.strip(), status_placeholder)

        st.session_state["results"] = results

# ── Display results ────────────────────────────────────────────────────────────
if "results" in st.session_state:
    res    = st.session_state["results"]
    report = res.get("final_report", {})

    # ── Timing strip ────────────────────────────────────────────────────────────
    t_cols = st.columns(5)
    times  = [
        ("Research",   res.get("research_time_s","?")),
        ("Merge",      res.get("merge_time_s","?")),
        ("Council",    res.get("council_time_s","?")),
        ("Synthesis",  res.get("synthesis_time_s","?")),
        ("Total",      res.get("total_time_s","?")),
    ]
    for col, (label, val) in zip(t_cols, times):
        col.metric(label, f"{val}s")

    st.markdown("---")

    # ── Tabs ─────────────────────────────────────────────────────────────────────
    tab_report, tab_council, tab_raw, tab_export = st.tabs([
        "📊 Final Report", "🏛️ Council Reviews", "🔍 Raw Research", "📁 Export"
    ])

    # ── Tab 1: Final Report ──────────────────────────────────────────────────────
    with tab_report:
        if report.get("parse_error"):
            st.error("Synthesis JSON parse error. Raw output shown below.")
            st.text(report.get("executive_summary",""))
        else:
            conf = report.get("overall_report_confidence", 0)
            conf_color = "#1D9E75" if conf >= 70 else "#BA7517" if conf >= 40 else "#E24B4A"

            # Confidence hero
            st.markdown(f"""
<div style="text-align:center; padding: 1.5rem; background:#F1EFE8; border-radius:12px; margin-bottom:1rem;">
  <div style="font-size:3.5rem; font-weight:700; color:{conf_color}">{conf}%</div>
  <div style="color:#5F5E5A; font-size:0.9rem; text-transform:uppercase; letter-spacing:1px">
    Overall Report Confidence
  </div>
</div>""", unsafe_allow_html=True)

            # Executive Summary
            st.markdown('<div class="section-header">Executive Summary</div>', unsafe_allow_html=True)
            st.info(report.get("executive_summary","N/A"))

            # Market Overview
            st.markdown('<div class="section-header">Market Overview</div>', unsafe_allow_html=True)
            mo = report.get("market_overview",{})
            m1, m2, m3 = st.columns(3)
            m1.markdown(f"""<div class="metric-card">
<div class="value">${mo.get('size_usd_billions','N/A')}B</div>
<div class="label">Market Size (USD)</div></div>""", unsafe_allow_html=True)
            m2.markdown(f"""<div class="metric-card">
<div class="value">{mo.get('cagr_percent','N/A')}%</div>
<div class="label">CAGR</div></div>""", unsafe_allow_html=True)
            m3.markdown(f"""<div class="metric-card">
<div class="value">{mo.get('confidence_score','N/A')}%</div>
<div class="label">Section Confidence</div></div>""", unsafe_allow_html=True)
            for pt in mo.get("key_data_points",[]):
                st.markdown(f"• {pt}")

            # Key Findings
            st.markdown('<div class="section-header">Key Findings</div>', unsafe_allow_html=True)
            for i, f in enumerate(report.get("key_findings",[]), 1):
                conf_f = f.get("confidence",0)
                bc = "#1D9E75" if conf_f>=70 else "#BA7517" if conf_f>=40 else "#E24B4A"
                st.markdown(f"""<div class="finding-row" style="border-left-color:{bc}">
<b>{i}.</b> {f.get('finding','')}
<span style="float:right; background:{bc}; color:white; padding:1px 8px;
border-radius:10px; font-size:0.75rem">{conf_f}%</span></div>""", unsafe_allow_html=True)

            # Competitive Landscape
            comp = report.get("competitive_landscape",[])
            if comp:
                st.markdown('<div class="section-header">Competitive Landscape</div>', unsafe_allow_html=True)
                import pandas as pd
                df_comp = pd.DataFrame(comp)
                st.dataframe(df_comp, use_container_width=True, hide_index=True)

            # Risks
            risks = report.get("risks",[])
            if risks:
                st.markdown('<div class="section-header">Risk Register</div>', unsafe_allow_html=True)
                for r in risks:
                    css = {"CRITICAL":"risk-crit","HIGH":"risk-high","MED":"risk-med","LOW":"risk-low"}
                    cls = css.get(r.get("severity",""),"risk-low")
                    st.markdown(f"""<div class="{cls}" style="padding:0.8rem; border-radius:8px; margin:0.4rem 0;">
<b>{r.get('severity','?')}</b> — {r.get('risk','')}
<br><small style="color:#5F5E5A">Mitigation: {r.get('mitigation','')}</small></div>""",
                        unsafe_allow_html=True)

            # Strategic Recommendations
            recs = sorted(report.get("strategic_recommendations",[]), key=lambda x: x.get("priority",9))
            if recs:
                st.markdown('<div class="section-header">Strategic Recommendations</div>', unsafe_allow_html=True)
                for rec in recs:
                    p    = rec.get("priority",3)
                    pclr = {1:"#1D9E75",2:"#BA7517",3:"#888780"}.get(p,"#888780")
                    st.markdown(f"""<div style="padding:0.7rem 1rem; margin:0.4rem 0; border-radius:8px;
background:#F8F8F6; border-left:4px solid {pclr}">
<span style="background:{pclr};color:white;padding:1px 8px;border-radius:10px;font-size:0.75rem;font-weight:700">P{p}</span>
&nbsp;<b>{rec.get('action','')}</b>
<span style="color:#BA7517; font-size:0.85rem"> — {rec.get('timeline','')}</span></div>""",
                        unsafe_allow_html=True)

            # Dissent log
            dissent = report.get("council_dissent_log",[])
            if dissent:
                with st.expander("⚠️ Council Dissent Log (unresolved disagreements)"):
                    for d in dissent:
                        st.markdown(f"**{d.get('agent','?')}**: {d.get('unresolved_challenge','')}")

    # ── Tab 2: Council Reviews ───────────────────────────────────────────────────
    with tab_council:
        reviews = res.get("council_reviews",[])
        if not reviews:
            st.info("No council reviews available.")
        else:
            for review in reviews:
                agent_name = review.get("agent","?")
                agent_obj  = next((a for a in COUNCIL_AGENTS if a.name == agent_name), None)
                clr        = agent_obj.color if agent_obj else "#888780"

                verdict     = review.get("verdict","?")
                verdict_bg  = {"ACCEPT":"#1D9E75","CHALLENGE":"#BA7517","REJECT":"#E24B4A"}.get(verdict,"#888780")
                adj         = review.get("confidence_adjustment",0)
                adj_str     = f"+{adj}" if adj > 0 else str(adj)

                with st.expander(f"**{agent_name}**  —  {review.get('focus','')}  |  Verdict: {verdict}  |  Confidence Δ: {adj_str}", expanded=True):
                    st.markdown(f"""
<span class="verdict-badge" style="background:{verdict_bg}">{verdict}</span>
<span style="margin-left:10px; font-size:0.85rem; color:#5F5E5A">
  Confidence adjustment: <b style="color:{'#1D9E75' if adj>0 else '#E24B4A'}">{adj_str}</b>
</span>""", unsafe_allow_html=True)

                    c1, c2 = st.columns(2)
                    with c1:
                        st.markdown("**✅ Validated Findings**")
                        for kf in review.get("key_findings",[]):
                            st.markdown(f"• {kf}")
                        st.markdown("**➕ Suggested Additions**")
                        for a in review.get("additions",[]):
                            st.markdown(f"• {a}")
                    with c2:
                        st.markdown("**⚠️ Challenges**")
                        for ch in review.get("challenges",[]):
                            st.markdown(f"• {ch}")
                        st.markdown("**🚩 Risk Flags**")
                        for rf in review.get("risk_flags",[]):
                            st.markdown(f"• {rf}")

    # ── Tab 3: Raw Research ──────────────────────────────────────────────────────
    with tab_raw:
        raw_tab1, raw_tab2, raw_tab3 = st.tabs(["Claude Output","Gemini Output","Merged Context"])
        with raw_tab1:
            st.markdown(res.get("claude_output","—"))
        with raw_tab2:
            st.markdown(res.get("gemini_output","—"))
        with raw_tab3:
            conflict_count = res.get("merged_context","").count("[CONFLICT:")
            if conflict_count:
                st.warning(f"Found {conflict_count} conflict(s) between Claude and Gemini outputs.")
            st.markdown(res.get("merged_context","—"))

    # ── Tab 4: Export ────────────────────────────────────────────────────────────
    with tab_export:
        st.markdown("### Download Reports")
        st.markdown(
            "Generate professional, boardroom-ready documents from the synthesised report. "
            "All exports include the full findings, risk register, and recommendations.")
        st.markdown("---")

        q_label  = res.get("query","report")[:40].replace(" ","_")
        now_str  = datetime.datetime.now().strftime("%Y%m%d")

        exp1, exp2, exp3 = st.columns(3)

        with exp1:
            st.markdown("#### 📄 PDF Report")
            st.markdown("Full A4 document with cover page, styled tables, risk colour-coding, and recommendations.")
            if st.button("Generate PDF", key="gen_pdf", use_container_width=True):
                if not REPORTLAB_AVAILABLE:
                    st.error("Install reportlab to enable PDF export.")
                else:
                    with st.spinner("Generating PDF..."):
                        try:
                            pdf_bytes = export_pdf(report, res["query"])
                            st.download_button(
                                "⬇️ Download PDF",
                                data=pdf_bytes,
                                file_name=f"market_research_{q_label}_{now_str}.pdf",
                                mime="application/pdf",
                                use_container_width=True,
                            )
                            st.success(f"PDF ready ({len(pdf_bytes)//1024} KB)")
                        except Exception as e:
                            st.error(f"PDF generation failed: {e}")

        with exp2:
            st.markdown("#### 📝 Word (DOCX)")
            st.markdown("Editable Word document with styled headings, colour-coded risk table, and numbered recommendations.")
            if st.button("Generate DOCX", key="gen_docx", use_container_width=True):
                if not DOCX_AVAILABLE:
                    st.error("Install python-docx to enable DOCX export.")
                else:
                    with st.spinner("Generating DOCX..."):
                        try:
                            docx_bytes = export_docx(report, res["query"])
                            st.download_button(
                                "⬇️ Download DOCX",
                                data=docx_bytes,
                                file_name=f"market_research_{q_label}_{now_str}.docx",
                                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                use_container_width=True,
                            )
                            st.success(f"DOCX ready ({len(docx_bytes)//1024} KB)")
                        except Exception as e:
                            st.error(f"DOCX generation failed: {e}")

        with exp3:
            st.markdown("#### 📊 PowerPoint (PPTX)")
            st.markdown("5-slide 16:9 deck: cover, exec summary with KPI cards, findings, risks, and recommendations.")
            if st.button("Generate PPTX", key="gen_pptx", use_container_width=True):
                if not PPTX_AVAILABLE:
                    st.error("Install python-pptx to enable PPTX export.")
                else:
                    with st.spinner("Generating PPTX..."):
                        try:
                            pptx_bytes = export_pptx(report, res["query"])
                            st.download_button(
                                "⬇️ Download PPTX",
                                data=pptx_bytes,
                                file_name=f"market_research_{q_label}_{now_str}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                use_container_width=True,
                            )
                            st.success(f"PPTX ready ({len(pptx_bytes)//1024} KB)")
                        except Exception as e:
                            st.error(f"PPTX generation failed: {e}")

        st.markdown("---")
        st.markdown("#### 📋 Raw JSON Report")
        st.markdown("Download the structured JSON for API integration or further processing.")
        json_bytes = json.dumps(report, indent=2).encode()
        st.download_button(
            "⬇️ Download JSON",
            data=json_bytes,
            file_name=f"market_research_{q_label}_{now_str}.json",
            mime="application/json",
        )

# ── Empty state ────────────────────────────────────────────────────────────────
else:
    st.markdown("""
<div style="text-align:center; padding:3rem; color:#888780">
  <div style="font-size:4rem">🔬</div>
  <div style="font-size:1.2rem; margin-top:1rem">Enter a research query and click <b>Run Research</b></div>
  <div style="margin-top:0.5rem; font-size:0.9rem">
    The pipeline will run Claude + Gemini in parallel, pass results through a 3-agent council,
    and synthesise a confidence-scored final report.
  </div>
</div>
""", unsafe_allow_html=True)
